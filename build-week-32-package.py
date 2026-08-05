from __future__ import annotations

import html
import json
import zipfile
from pathlib import Path

from PIL import Image, ImageStat
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent
WEEK = "2026-week-32"
SPEC_DIR = ROOT / "specs" / WEEK
CHART_DIR = ROOT / "charts" / WEEK

ORDER = [
    "ozon-insurance-risk",
    "marketplace-retrenchment",
    "wildberries-support-request",
    "wildberries-kazakhstan-capacity",
    "wildberries-warehouse-damage-map",
    "ecommerce-company-churn",
    "ecommerce-active-company-base",
    "ecommerce-growth-deceleration",
    "russia-fuel-regional-map",
    "zabaykalsky-fuel-coverage",
    "fuel-export-ban-structure",
    "emergency-fuel-imports",
    "refinery-output-gap",
    "mongolia-fuel-buffer",
    "vtb-h1-profit",
    "vtb-other-operating-income",
    "vtb-loan-book",
    "grain-domestic-prices",
    "grain-export-prices",
    "azov-grain-export-loss",
    "black-sea-terminal-exposure",
    "sunflower-export-shock",
    "east-coal-demand-slump",
    "turkey-coal-squeeze",
    "urals-discount-freight",
    "south-moscow-sublease",
    "gtlk-first-half-loss",
    "china-tire-pressure",
    "fashion-volume-value-split",
    "alrosa-profit-reversal",
    "russian-inflation-dashboard",
]


def read_spec(slug: str) -> dict:
    return json.loads((SPEC_DIR / f"{slug}.json").read_text(encoding="utf-8"))


def corner_background(image: Image.Image) -> tuple[int, int, int]:
    rgb = image.convert("RGB")
    w, h = rgb.size
    sample = max(8, min(w, h) // 40)
    boxes = [
        (0, 0, sample, sample),
        (w - sample, 0, w, sample),
        (0, h - sample, sample, h),
        (w - sample, h - sample, w, h),
    ]
    crops = [rgb.crop(box) for box in boxes]
    means = [ImageStat.Stat(crop).mean for crop in crops]
    return tuple(round(sum(mean[i] for mean in means) / len(means)) for i in range(3))


def build_pptx() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(8.333333)
    prs.core_properties.title = "Stanichart Data Stories — 2026 Week 32"
    prs.core_properties.subject = "31 branded data-story charts"
    prs.core_properties.author = "Stanichart"
    prs.core_properties.comments = "One chart per slide, generated from the validated ChartSpecs."

    blank = prs.slide_layouts[6]
    margin = Inches(0.04)
    max_w = prs.slide_width - 2 * margin
    max_h = prs.slide_height - 2 * margin

    for slug in ORDER:
        png = CHART_DIR / f"{slug}.png"
        spec = read_spec(slug)
        if not png.exists():
            raise FileNotFoundError(png)

        with Image.open(png) as image:
            width_px, height_px = image.size
            bg = corner_background(image)

        slide = prs.slides.add_slide(blank)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*bg)

        image_ratio = width_px / height_px
        frame_ratio = max_w / max_h
        if image_ratio >= frame_ratio:
            width = max_w
            height = int(width / image_ratio)
        else:
            height = max_h
            width = int(height * image_ratio)
        left = int((prs.slide_width - width) / 2)
        top = int((prs.slide_height - height) / 2)
        picture = slide.shapes.add_picture(str(png), left, top, width=width, height=height)
        picture._element.nvPicPr.cNvPr.set("descr", spec["title"])

    output = CHART_DIR / "stanichart-week-32.pptx"
    prs.save(output)
    return output


def write_audit() -> Path:
    rows = []
    for position, slug in enumerate(ORDER, start=1):
        png = CHART_DIR / f"{slug}.png"
        html_file = CHART_DIR / f"{slug}.html"
        spec = read_spec(slug)
        with Image.open(png) as image:
            rows.append(
                {
                    "slide": position,
                    "slug": slug,
                    "title": spec["title"],
                    "recipe": spec["recipe"],
                    "png": png.name,
                    "html": html_file.name,
                    "width": image.width,
                    "height": image.height,
                    "mode": image.mode,
                    "pngBytes": png.stat().st_size,
                    "htmlBytes": html_file.stat().st_size,
                    "backgroundRgb": list(corner_background(image)),
                }
            )
    audit = {
        "week": WEEK,
        "chartCount": len(rows),
        "allPngsPresent": all((CHART_DIR / f"{slug}.png").exists() for slug in ORDER),
        "allHtmlPresent": all((CHART_DIR / f"{slug}.html").exists() for slug in ORDER),
        "charts": rows,
    }
    output = CHART_DIR / "image-audit.json"
    output.write_text(json.dumps(audit, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    (CHART_DIR / "slide-order.json").write_text(
        json.dumps({"week": WEEK, "slides": rows}, indent=2, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )
    return output


def write_gallery() -> Path:
    cards = []
    for index, slug in enumerate(ORDER, start=1):
        spec = read_spec(slug)
        title = html.escape(spec["title"])
        subtitle = html.escape(spec.get("subtitle", ""))
        recipe = html.escape(spec["recipe"])
        cards.append(
            f'''<article class="card">
  <a href="{slug}.html"><img src="{slug}.png" alt="{title}" loading="lazy"></a>
  <div class="meta"><span>{index:02d}</span><span>{recipe}</span></div>
  <h2><a href="{slug}.html">{title}</a></h2>
  <p>{subtitle}</p>
</article>'''
        )
    document = f'''<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>Stanichart — {WEEK}</title>
<style>
:root {{ color-scheme: dark; --bg:#101519; --panel:#182026; --text:#f1f3ef; --muted:#aab4b8; --line:#354148; --accent:#f4c34f; }}
* {{ box-sizing:border-box; }}
body {{ margin:0; background:var(--bg); color:var(--text); font-family:Arial, Helvetica, sans-serif; }}
header {{ max-width:1500px; margin:auto; padding:48px 28px 28px; }}
h1 {{ margin:0; font-size:clamp(34px,5vw,68px); line-height:.98; letter-spacing:-.04em; }}
header p {{ color:var(--muted); max-width:760px; font-size:18px; line-height:1.5; }}
.grid {{ max-width:1500px; margin:auto; padding:16px 28px 72px; display:grid; grid-template-columns:repeat(auto-fit,minmax(360px,1fr)); gap:24px; }}
.card {{ background:var(--panel); border:1px solid var(--line); border-radius:16px; overflow:hidden; box-shadow:0 10px 26px rgba(0,0,0,.22); }}
.card img {{ width:100%; aspect-ratio:1.2; object-fit:contain; background:#f7f6f1; display:block; }}
.card .meta {{ display:flex; justify-content:space-between; padding:20px 22px 0; color:var(--accent); text-transform:uppercase; letter-spacing:.08em; font-size:12px; font-weight:700; }}
.card h2 {{ margin:10px 22px 8px; font-size:24px; line-height:1.12; letter-spacing:-.02em; }}
.card h2 a {{ color:var(--text); text-decoration:none; }}
.card p {{ margin:0 22px 24px; color:var(--muted); line-height:1.45; }}
@media (max-width:520px) {{ .grid {{ grid-template-columns:1fr; padding-left:14px; padding-right:14px; }} header {{ padding-left:18px; padding-right:18px; }} }}
</style>
</head>
<body>
<header>
  <h1>Stanichart<br>{WEEK}</h1>
  <p>Thirty-one branded data stories, each available as responsive HTML and a publication-ready PNG. Select a chart to open the interactive-quality full layout.</p>
</header>
<main class="grid">{''.join(cards)}</main>
</body>
</html>
'''
    output = CHART_DIR / "index.html"
    output.write_text(document, encoding="utf-8")
    return output


def write_readme(pptx_path: Path) -> Path:
    lines = [
        f"# Stanichart chart package — {WEEK}",
        "",
        "This folder contains 31 branded charts, one for each data story identified in `input.txt`.",
        "",
        "## Primary deliverables",
        "",
        "- `index.html` — visual gallery linking every responsive chart.",
        "- `stanichart-week-32.pptx` — 31-slide deck, one chart per slide.",
        "- `*.png` — publication-ready chart captures.",
        "- `*.html` — responsive branded chart pages.",
        "- `image-audit.json` — image dimensions, file sizes and slide mapping.",
        "- `slide-order.json` — editorial slide order.",
        "",
        "## Coverage and routing",
        "",
        "- `../../specs/2026-week-32/coverage-matrix.md` maps every input datapoint to a chart.",
        "- `../../specs/2026-week-32/routing-matrix.md` records the standard-chart versus regional-map decisions.",
        "",
        f"Deck size: {pptx_path.stat().st_size:,} bytes.",
    ]
    output = CHART_DIR / "README.md"
    output.write_text("\n".join(lines) + "\n", encoding="utf-8")
    return output


def build_zip(pptx_path: Path) -> Path:
    output = CHART_DIR / "stanichart-week-32-package.zip"
    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9) as archive:
        for slug in ORDER:
            for suffix in (".png", ".html"):
                path = CHART_DIR / f"{slug}{suffix}"
                archive.write(path, f"charts/{path.name}")
            spec = SPEC_DIR / f"{slug}.json"
            archive.write(spec, f"specs/{spec.name}")
        for path in (
            pptx_path,
            CHART_DIR / "index.html",
            CHART_DIR / "README.md",
            CHART_DIR / "image-audit.json",
            CHART_DIR / "slide-order.json",
            SPEC_DIR / "manifest.json",
            SPEC_DIR / "coverage-matrix.md",
            SPEC_DIR / "routing-matrix.md",
        ):
            archive.write(path, path.name)
    return output


def main() -> None:
    if len(ORDER) != 31 or len(set(ORDER)) != 31:
        raise RuntimeError("Expected exactly 31 unique chart slugs")
    expected_specs = {p.stem for p in SPEC_DIR.glob("*.json") if p.name != "manifest.json"}
    if expected_specs != set(ORDER):
        missing = sorted(set(ORDER) - expected_specs)
        extra = sorted(expected_specs - set(ORDER))
        raise RuntimeError(f"Spec/order mismatch; missing={missing}, extra={extra}")

    audit_path = write_audit()
    gallery_path = write_gallery()
    pptx_path = build_pptx()
    readme_path = write_readme(pptx_path)
    zip_path = build_zip(pptx_path)

    print(
        json.dumps(
            {
                "charts": len(ORDER),
                "audit": str(audit_path.relative_to(ROOT)),
                "gallery": str(gallery_path.relative_to(ROOT)),
                "pptx": str(pptx_path.relative_to(ROOT)),
                "readme": str(readme_path.relative_to(ROOT)),
                "zip": str(zip_path.relative_to(ROOT)),
                "pptxBytes": pptx_path.stat().st_size,
                "zipBytes": zip_path.stat().st_size,
            },
            indent=2,
        )
    )


if __name__ == "__main__":
    main()
